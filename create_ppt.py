# -*- coding: utf-8 -*-
"""
生成参赛路演PPT
运行: python create_ppt.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# 配色
PRI = RGBColor(0x1A, 0x56, 0xDB)
GRN = RGBColor(0x10, 0xB9, 0x81)
RED = RGBColor(0xEF, 0x44, 0x44)
DRK = RGBColor(0x1F, 0x29, 0x37)
LGT = RGBColor(0xF3, 0xF4, 0xF6)
WHT = RGBColor(0xFF, 0xFF, 0xFF)
ORG = RGBColor(0xF5, 0x9E, 0x0B)
PUR = RGBColor(0x8B, 0x5C, 0xF6)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def bg(slide, color=LGT):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color

def bar(slide, text, sub=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.3))
    s.fill.solid(); s.fill.fore_color.rgb = PRI; s.line.fill.background()
    tf = s.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text; p.font.size = Pt(36)
    p.font.color.rgb = WHT; p.font.bold = True; p.alignment = PP_ALIGN.LEFT
    tf.margin_left = Inches(0.8); tf.margin_top = Inches(0.25)
    if sub:
        ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(1.3), prs.slide_width, Inches(0.06))
        ln.fill.solid(); ln.fill.fore_color.rgb = GRN; ln.line.fill.background()

def tb(slide, l, t, w, h, text, fs=16, color=DRK, bold=False, align=PP_ALIGN.LEFT):
    bx = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = bx.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text; p.font.size = Pt(fs)
    p.font.color.rgb = color; p.font.bold = bold; p.alignment = align
    return tf

def card(slide, l, t, w, h, title, lines, tc=PRI):
    c = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    c.fill.solid(); c.fill.fore_color.rgb = WHT
    c.line.color.rgb = RGBColor(0xE5, 0xE7, 0xEB); c.line.width = Pt(1)
    tf = c.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.2); tf.margin_right = Inches(0.2); tf.margin_top = Inches(0.15)
    p = tf.paragraphs[0]; p.text = title; p.font.size = Pt(18)
    p.font.color.rgb = tc; p.font.bold = True
    for ln in lines:
        p = tf.add_paragraph(); p.text = ln; p.font.size = Pt(13)
        p.font.color.rgb = DRK; p.space_before = Pt(3)

def bbar(slide, text='因果增强工业智能体 · 创智青山AI智能体创新大赛 · 技术挑战赛道'):
    b = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(7.0), prs.slide_width, Inches(0.5))
    b.fill.solid(); b.fill.fore_color.rgb = DRK; b.line.fill.background()
    tf = b.text_frame; tf.margin_top = Inches(0.08)
    p = tf.paragraphs[0]; p.text = text; p.font.size = Pt(11)
    p.font.color.rgb = RGBColor(0x9C, 0xA3, 0xAF); p.alignment = PP_ALIGN.CENTER

# ================================================================
# S1: 封面
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, PRI)
tb(s, 1, 1.5, 11.3, 1.5, '因果增强工业智能体', 52, WHT, True, PP_ALIGN.CENTER)
tb(s, 1, 2.7, 11.3, 1, 'Causal Enhanced Industrial Agent', 28, RGBColor(0x93, 0xC5, 0xFD), align=PP_ALIGN.CENTER)
ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.5), Inches(3.6), Inches(4.3), Inches(0.04))
ln.fill.solid(); ln.fill.fore_color.rgb = GRN; ln.line.fill.background()
tb(s, 1, 4.0, 11.3, 0.8, '基于双通道因果图融合的工业异常根因分析与反事实推理系统', 20, RGBColor(0xD1, 0xD5, 0xDB), align=PP_ALIGN.CENTER)
tb(s, 1, 5.5, 11.3, 0.5, '创智青山AI智能体创新大赛 · 技术挑战赛道 · 钢铁石化新材料方向', 16, WHT, align=PP_ALIGN.CENTER)

# ================================================================
# S2: 产业痛点
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); bar(s, '产业痛点：传统工业AI为什么不够？')
pain = [
    ('🔴 痛点1：只会报警，不会解释', [
        '传统异常检测(PCA/OCSVM)告诉你"反应器温度高了"',
        '但不知道为什么高 — 是冷却水问题？进料问题？还是传感器坏了？',
        '操作员面对大量报警不知所措，容易误判根因导致错误操作']),
    ('🟠 痛点2：各系统各自为政，缺乏全局视角', [
        '巡检机器人、DCS报警、振动监测、油液分析各干各的',
        '没有人把"阀门卡滞→冷却水不足→温度升高"这条链串起来',
        '结果是头痛医头、脚痛医脚']),
    ('🔵 痛点3：大模型直接推理有幻觉风险', [
        '大模型看到"温度高了"会给出看似合理但可能错误的解释',
        '工业场景零容错 — 一次误判可能造成数百万损失',
        '需要物理因果约束来抑制幻觉']),
]
for i, (ti, ln) in enumerate(pain):
    card(s, 0.5 + i*4.2, 1.8, 3.9, 3.8, ti, ln)
tb(s, 0.8, 6.0, 11.7, 0.8, '核心命题：如何让AI不仅感知"出了什么问题"，还能理解"为什么会出问题"并给出"应该怎么做"？', 17, RED, True, PP_ALIGN.CENTER)
bbar(s)

# ================================================================
# S3: 技术方案总览
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); bar(s, '技术方案：从"感知"到"理解"到"决策"的因果闭环')
steps = [
    ('📡 感知层', '异常检测', ['从正常工况数据自动学习基线', '多变量同时监控', '输出：哪些变量偏离了正常范围'], PRI),
    ('🧠 理解层', '因果根因分析', ['在因果图上逆流回溯', '双通道融合：知识+数据交叉验证', '输出：根因排序 + 因果路径 + 证据链'], GRN),
    ('💡 决策层', '反事实推理', ['"如果修复阀门，温度能恢复多少？"', '对比多种干预方案的效果', '输出：推荐操作 + 量化预期改善'], PUR),
]
for i, (ly, ti, ln, cl) in enumerate(steps):
    lt = 0.5 + i*4.2
    cr = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(lt+1.3), Inches(1.7), Inches(1.0), Inches(1.0))
    cr.fill.solid(); cr.fill.fore_color.rgb = cl; cr.line.fill.background()
    tf = cr.text_frame; tf.margin_top = Inches(0.15)
    p = tf.paragraphs[0]; p.text = ly; p.font.size = Pt(20); p.font.color.rgb = WHT; p.alignment = PP_ALIGN.CENTER
    if i < 2:
        ar = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(lt+3.0), Inches(2.0), Inches(0.8), Inches(0.4))
        ar.fill.solid(); ar.fill.fore_color.rgb = DRK; ar.line.fill.background()
    card(s, lt, 3.2, 3.9, 3.3, ti, ln, cl)
tb(s, 0.8, 6.2, 11.7, 0.6, '传统AI：数据→模型→"温度异常" ← 结束     |     本方案：数据→因果图→"温度异常，根因：阀门卡滞，建议：恢复开度" ← 闭环', 15, RED, True, PP_ALIGN.CENTER)
bbar(s)

# ================================================================
# S4: 核心创新1 — 双通道因果图融合
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); bar(s, '核心创新①：双通道因果图融合算法', '知识驱动 + 数据驱动 → 交叉验证 → 鲁棒因果图')
card(s, 0.5, 1.8, 5.8, 2.5, '📖 通道1：知识驱动', [
    '输入：操作手册、维修规程等非结构化文档',
    '方法：LLM抽取因果对 + 规则匹配fallback',
    '示例：CW_Valve → CW_Flow (阀门开度→冷却水流量)',
    '优势：覆盖慢因果、稀有事件、专家经验',
], PRI)
card(s, 6.8, 1.8, 5.8, 2.5, '📊 通道2：数据驱动', [
    '输入：传感器历史时序数据',
    '方法：PCMCI+ 条件独立性检验',
    '示例：CW_Flow → Reactor_Temp (p<0.05)',
    '优势：客观、发现文档未记录的隐含因果',
], GRN)
card(s, 2.0, 4.8, 9.3, 2.0, '🔄 融合层：冲突消解 + 置信度加权 ⭐ 原创算法', [
    '知识有+数据有 → 双重验证，最高置信度  |  知识有+数据无 → 降权，标注"待数据验证"',
    '知识无+数据有 → 标注"数据驱动新发现"  |  均无 → 不建边',
    '创新：不是简单取并集，而是利用两通道互补优势（知识覆盖慢因果、数据覆盖快因果）',
], ORG)
bbar(s)

# ================================================================
# S5: 核心创新2 — 根因分析对比
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); bar(s, '核心创新②：因果根因分析 — 从"发现异常"到"理解异常"')
card(s, 0.5, 1.8, 5.8, 4.5, '❌ 传统异常检测', [
    'Reactor_Temp = 192°C ⚠ 偏高',
    'CW_Flow = 118 ⚠ 偏低, CW_Valve = 22% ⚠ 偏低',
    'Product_Conc = 0.57 ⚠ 偏低',
    '——',
    '→ 操作员: 8个报警同时响，哪个是原因？',
    '→ 传统系统: 给不出答案，靠人判断',
], RED)
card(s, 6.8, 1.8, 5.8, 4.5, '✅ 因果根因分析（本方案）', [
    'Reactor_Temp 异常 ↑',
    '  └→ 根因#1: CW_Flow (评分0.74)',
    '       └→ 根因: CW_Valve (评分0.65)',
    '           因果链: CW_Valve→CW_Flow→Reactor_Temp',
    '',
    '处置建议: 检查阀门卡滞 → 切换备用回路',
    '→ 因果链每步都有证据，可追溯可解释',
], GRN)
tb(s, 0.8, 5.5, 11.7, 1.2, '算法：①检测异常 → ②因果图获取祖先 → ③筛选祖先中同样异常的 → ④按因果效应×置信度×路径距离×偏离程度综合评分排序', 14, DRK)
bbar(s)

# ================================================================
# S6: 核心创新3 — 反事实推理
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); bar(s, '核心创新③：反事实推理 — "如果...会怎样？"')
card(s, 0.5, 1.8, 4.0, 4.5, '🤔 反事实问题', [
    '如果修复冷却水阀门(开度→60%),',
    '反应器温度会降到多少？',
    '',
    '如果同时修复阀门+降低进料流量,',
    '产物浓度能恢复正常吗？',
    '',
    '哪个干预方案效果最好？',
], PRI)
card(s, 5.0, 1.8, 4.0, 4.5, '⚙️ 推理机制', [
    '三步反事实推理框架：',
    '① 溯因(Abduction): 从观测推断噪声',
    '② 行动(Action): 施加干预 do(X=x)',
    '③ 预测(Prediction): 沿因果图传播效应',
    '',
    '基于 Pearl 结构因果模型(SCM)',
    '线性SEM解析可解 / 含环图迭代收敛',
], GRN)
card(s, 9.5, 1.8, 3.3, 4.5, '📊 推理结果', [
    '方案1: 修复阀门',
    '  温度: 192→165°C (-27)',
    '方案2: 降低进料',
    '  温度: 192→180°C (-12)',
    '方案3: 综合干预 ✓',
    '  温度: 192→155°C (-37)',
    '归因: 阀门58% 进料42%',
], PUR)
bbar(s)

# ================================================================
# S7: LLM角色
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); bar(s, 'LLM的角色：建图工具，不是推理工具')
card(s, 0.5, 1.8, 5.8, 3.0, '✅ LLM做的事（仅一次，初始化时）', [
    '输入：操作手册、维修规程等非结构化文档',
    '输出："CW_Valve→CW_Flow (阀门开度→冷却水流量)"',
    '本质：把工程师写的文字翻译成结构化的因果对JSON',
    '频率：离线运行一次，结果存下来后续复用',
], GRN)
card(s, 6.8, 1.8, 5.8, 3.0, '❌ LLM不做的事（运行时，绝不调用）', [
    '不做：看到"Reactor_Temp=192"→回答"可能是阀门问题"',
    '不做：实时异常判断 ← 统计方法的事',
    '不做：根因路径推理 ← 图搜索算法的事',
    '不做：反事实计算 ← 结构因果模型的代数计算',
], RED)
tb(s, 0.8, 5.2, 11.7, 1.5, 'LLM擅长理解文本→用来"建图"。因果图擅长推理→用来"诊断"。各取所长，互不越界。\n答辩关键：如果评委问"为什么不直接用大模型诊断？"→ 回答：工业零容错，因果图提供物理约束，从根本上抑制幻觉。', 14, DRK)
bbar(s)

# ================================================================
# S8: 实验结果
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); bar(s, '实验验证：合成数据 + 真实TEP工业数据双重验证')
card(s, 0.5, 1.8, 6.0, 2.8, '🔬 合成CSTR数据验证（算法正确性）', [
    '场景：12变量CSTR反应器，4种故障模式，已知Ground Truth(17条边)',
    '根因准确率：3/4 (阀门卡滞✓ 入口温度高✓ 进料突增✓)',
    '双通道融合：知识15条+数据6条→融合21条',
    '反事实推理：正确评估3种干预方案并推荐最优',
], PRI)
card(s, 6.8, 1.8, 6.0, 2.8, '🏭 真实TEP数据验证（实用性）', [
    '数据集：Tennessee Eastman Process (MIT Braatz Group)',
    '规模：52过程变量，21种故障，每故障~1440采样点',
    '',
    '✓ IDV(6) A进料损失: 因果链准确检出',
    '✓ IDV(14) 阀门卡滞: Valve_Reactor_CW→Reactor_Temp 根因正确',
], GRN)
card(s, 0.5, 5.0, 12.3, 1.7, '📊 数据与工具', [
    '合成数据：自研CSTR仿真器(12变量,4故障) | TEP数据：MIT Braatz Group (4.7MB,公开数据集)',
    '因果发现：PCMCI+ (Tigramite v5) | 因果推断：Pearl SCM框架 | LLM：Claude API(可选)',
    '前端：Streamlit+Plotly 交互可视化 | 环境：Python 3.11,纯CPU可跑,单步推理<10ms',
], DRK)
bbar(s)

# ================================================================
# S9: 方法对比
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); bar(s, '与现有方法的本质差异')
# 简化的对比表用卡片形式
comp = [
    ('传统异常检测\n(PCA/OCSVM)', '✓ 异常检测\n✗ 根因分析\n✗ 反事实\n✗ 处置建议\n可解释性: 低', RED),
    ('纯大模型方案\n(LLM直接诊断)', '✓ 异常检测\n△ 根因(有幻觉)\n✗ 反事实\n△ 处置(不保证对)\n可解释性: 中', ORG),
    ('纯数据因果发现\n(仅PCMCI+)', '✓ 异常检测\n△ 根因(无解释)\n△ 反事实(仅线性)\n✗ 处置建议\n可解释性: 中', ORG),
    ('本方案\n(因果增强智能体)', '✓ 异常检测\n✓ 根因(双通道验证)\n✓ 反事实推理\n✓ 处置建议(因果约束)\n可解释性: 高', GRN),
]
for i, (ti, ct, cl) in enumerate(comp):
    card(s, 0.5 + i*3.2, 1.8, 2.9, 3.2, ti, ct.split('\n'), cl)
tb(s, 0.8, 5.5, 11.7, 1.2, '本质差异：范式升级，不是精度优化。传统停留在感知层，大模型有幻觉，本方案实现从感知→理解→决策的完整因果闭环。因果图提供物理约束→推理每一步可追溯、可验证、可解释。', 14, RED, True)
bbar(s)

# ================================================================
# S10: 落地路径
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); bar(s, '落地路径：对接青山产业，从Demo到产线')
path_steps = [
    ('Phase 1\n（当前）', '原型验证', ['✓ 合成数据验证', '✓ 真实TEP验证', '✓ Streamlit Demo', '→ 参赛答辩'], PRI),
    ('Phase 2\n（3个月）', '企业试点', ['对接武钢/石化企业', '获取工艺文档+数据', '构建企业因果图', '离线诊断测试'], GRN),
    ('Phase 3\n（6个月）', '系统集成', ['对接PLC/SCADA', '实时数据流推理', '接入报警工单系统', '操作员培训'], PUR),
    ('Phase 4\n（12个月）', '规模化', ['推广多产线/企业', '行业因果知识库', 'SaaS化部署', '持续优化'], ORG),
]
for i, (ph, ti, ln, cl) in enumerate(path_steps):
    card(s, 0.5 + i*3.2, 1.8, 2.9, 4.8, f'{ph}\n{ti}', ln, cl)
bbar(s)

# ================================================================
# S11: 总结
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, PRI)
tb(s, 1, 0.8, 11.3, 1, '总结', 40, WHT, True, PP_ALIGN.CENTER)
items = [
    ('🎯', '核心命题', '让工业AI从"能看见"走到"能理解、能决策"'),
    ('⚙️', '技术路径', '双通道因果图融合 + 因果根因分析 + 反事实推理'),
    ('💡', '关键创新', '不是精度改进，是能力维度升级 — 范式级创新'),
    ('🔒', '核心壁垒', '因果图提供物理约束，推理每步可追溯、可验证'),
    ('🏭', '落地价值', '直接对接青山钢铁石化产业，从故障诊断到工艺优化'),
]
for i, (ic, ti, ds) in enumerate(items):
    y = 2.0 + i*0.9
    tb(s, 1.5, y, 0.6, 0.6, ic, 28, WHT, align=PP_ALIGN.CENTER)
    color = GRN if i==4 else RGBColor(0xD1, 0xD5, 0xDB)
    tb(s, 2.5, y, 9.8, 0.8, f'{ti}: {ds}', 18, color)
ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2.0), Inches(6.3), Inches(9.3), Inches(0.04))
ln.fill.solid(); ln.fill.fore_color.rgb = GRN; ln.line.fill.background()
tb(s, 1, 6.5, 11.3, 0.5, '一句话：我们不造另一个只会报警的AI，我们造一个能告诉操作员"为什么"和"怎么办"的AI工程师。', 20, WHT, True, PP_ALIGN.CENTER)

# ================================================================
out = 'f:/创新/causal-industrial-agent/presentation.pptx'
prs.save(out)
print(f'✅ PPT已生成: {out}')
print(f'共 {len(prs.slides)} 页')
